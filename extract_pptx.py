import sys
from pptx import Presentation

pptx_path = r"C:\Users\UserPC\Desktop\안티그래비티\교육플렛폼\중부위험성평가추가분.pptx"
prs = Presentation(pptx_path)

texts = []

def extract_text(shape):
    try:
        if shape.has_text_frame:
            texts.append(shape.text)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if hasattr(cell, 'text_frame') and cell.text_frame:
                        texts.append(cell.text)
        if shape.shape_type == 6: # GROUP
            for s in shape.shapes:
                extract_text(s)
    except Exception:
        pass

for slide in prs.slides:
    for shape in slide.shapes:
        extract_text(shape)

with open(r"C:\Users\UserPC\Desktop\안티그래비티\교육플렛폼\extracted.txt", "w", encoding="utf-8") as f:
    f.write("\n".join(texts))
print("Extraction complete.")
